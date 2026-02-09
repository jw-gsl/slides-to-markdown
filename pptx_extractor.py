#!/usr/bin/env python3

from pptx import Presentation
import json
import platform
import shutil
import argparse
from pathlib import Path

SCRIPT_DIR = Path(__file__).resolve().parent
CONFIG_FILE = SCRIPT_DIR / '.slides-config.json'


def extract_table_as_markdown(table):
    """Convert a PowerPoint table to a markdown table."""
    rows = []
    for row in table.rows:
        cells = [cell.text.strip().replace('|', '\\|') for cell in row.cells]
        rows.append(cells)

    if not rows:
        return ""

    col_count = max(len(row) for row in rows)
    # Pad rows to same length
    rows = [row + [''] * (col_count - len(row)) for row in rows]

    header = '| ' + ' | '.join(rows[0]) + ' |'
    separator = '| ' + ' | '.join(['---'] * col_count) + ' |'
    body_rows = ['| ' + ' | '.join(row) + ' |' for row in rows[1:]]

    return '\n'.join([header, separator] + body_rows)


def extract_text_from_shape(shape):
    """Extract text from a shape with markdown formatting."""
    text_items = []

    # Tables — render as markdown tables (skip shape.text to avoid duplicates)
    if shape.has_table:
        table_md = extract_table_as_markdown(shape.table)
        if table_md:
            text_items.append(table_md)
        return text_items

    # Text frames — handle paragraph levels for bullet formatting
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            text = paragraph.text.strip()
            if text:
                level = paragraph.level
                if level > 0:
                    indent = '  ' * (level - 1)
                    text_items.append(f"{indent}- {text}")
                else:
                    text_items.append(text)
        return text_items

    # Grouped shapes
    if shape.shape_type == 6:  # GROUP shape type
        for sub_shape in shape.shapes:
            text_items.extend(extract_text_from_shape(sub_shape))

    return text_items


def extract_notes(slide):
    """Extract speaker notes from a slide."""
    try:
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                return notes_text
    except Exception:
        pass
    return None


def extract_text_from_pptx(file_path):
    """Extract all text from a PPTX file and return as markdown."""
    prs = Presentation(file_path)
    text_runs = []

    # Add document title as h1
    base_name = Path(file_path).stem
    text_runs.append(f"# {base_name}\n")

    for slide_number, slide in enumerate(prs.slides, start=1):
        # Use slide title as heading if available
        title_shape = slide.shapes.title
        if title_shape and title_shape.text.strip():
            text_runs.append(f"## {title_shape.text.strip()}\n")
        else:
            text_runs.append(f"## Slide {slide_number}\n")

        for shape in slide.shapes:
            # Skip the title shape — already used as the heading
            if title_shape is not None and shape.shape_id == title_shape.shape_id:
                continue
            text_runs.extend(extract_text_from_shape(shape))

        # Append speaker notes as blockquote
        notes = extract_notes(slide)
        if notes:
            quoted = '\n'.join(f'> {line}' for line in notes.split('\n'))
            text_runs.append(f"\n**Notes:**\n{quoted}")

    return "\n\n".join(text_runs)


def create_folder_structure(base_dir, quiet=False):
    """Create the necessary folder structure."""
    folders = {
        'input': Path(base_dir) / 'input',
        'processed': Path(base_dir) / 'processed',
        'output': Path(base_dir) / 'output'
    }

    for folder_name, folder_path in folders.items():
        if not folder_path.exists():
            folder_path.mkdir(parents=True, exist_ok=True)
            if not quiet:
                print(f"Created folder: {folder_path}")
        elif not quiet:
            print(f"Verified folder: {folder_path}")

    return folders


# --- Config & interactive setup ---

def load_config():
    """Load saved config, or return None if no config exists."""
    if CONFIG_FILE.exists():
        try:
            with open(CONFIG_FILE, 'r') as f:
                return json.load(f)
        except (json.JSONDecodeError, OSError):
            pass
    return None


def save_config(working_dir):
    """Save the chosen working directory to config."""
    config = {'working_dir': str(working_dir)}
    with open(CONFIG_FILE, 'w') as f:
        json.dump(config, f, indent=2)
    print(f"\nSaved to {CONFIG_FILE}")
    print(f"Future runs will use: {working_dir}")
    print(f"(Override anytime with --dir or re-run --setup)\n")


def get_default_options():
    """Return a list of sensible default directory options for the current platform."""
    home = Path.home()
    options = []

    # Option 1: script directory (always available)
    options.append(('Here (next to the script)', SCRIPT_DIR))

    # Platform-specific suggestions
    if platform.system() == 'Darwin':
        desktop = home / 'Desktop' / 'slides-to-markdown'
        documents = home / 'Documents' / 'slides-to-markdown'
        options.append(('Desktop', desktop))
        options.append(('Documents', documents))
    elif platform.system() == 'Windows':
        desktop = home / 'Desktop' / 'slides-to-markdown'
        documents = home / 'Documents' / 'slides-to-markdown'
        options.append(('Desktop', desktop))
        options.append(('Documents', documents))
    else:
        options.append(('Home folder', home / 'slides-to-markdown'))

    return options


def interactive_setup():
    """Guide the user through choosing a working directory."""
    print("=" * 50)
    print("  Slides to Markdown — First-time Setup")
    print("=" * 50)
    print()
    print("Where would you like the working folders")
    print("(input/, output/, processed/) to be created?")
    print()

    options = get_default_options()

    for i, (label, path) in enumerate(options, 1):
        print(f"  [{i}] {label}")
        print(f"      {path}")
        print()

    custom_num = len(options) + 1
    print(f"  [{custom_num}] Custom path")
    print()

    while True:
        try:
            choice = input(f"Choose an option [1-{custom_num}]: ").strip()
            choice_num = int(choice)

            if 1 <= choice_num <= len(options):
                _, chosen_path = options[choice_num - 1]
                break
            elif choice_num == custom_num:
                custom = input("Enter the full path: ").strip()
                if not custom:
                    print("No path entered. Please try again.")
                    continue
                chosen_path = Path(custom).expanduser().resolve()
                break
            else:
                print(f"Please enter a number between 1 and {custom_num}.")
        except ValueError:
            print(f"Please enter a number between 1 and {custom_num}.")
        except (KeyboardInterrupt, EOFError):
            print("\nSetup cancelled.")
            return None

    chosen_path = Path(chosen_path).resolve()

    print(f"\nSetting up in: {chosen_path}")
    create_folder_structure(chosen_path)
    save_config(chosen_path)

    return chosen_path


def resolve_working_dir(cli_dir=None):
    """Determine the working directory from CLI args, saved config, or interactive setup."""
    # Explicit --dir always wins
    if cli_dir and cli_dir != '.':
        return Path(cli_dir).resolve()

    # Check for saved config
    config = load_config()
    if config and 'working_dir' in config:
        saved = Path(config['working_dir'])
        if saved.exists():
            return saved.resolve()

    # No config and no explicit dir — default to script directory
    return SCRIPT_DIR


def process_pptx_files(base_dir, keep=False):
    """Process all PPTX files in the input directory."""
    folders = create_folder_structure(base_dir, quiet=True)

    input_dir = folders['input']
    processed_dir = folders['processed']
    output_dir = folders['output']

    # Find all PPTX files in input directory (case-insensitive)
    pptx_files = list(input_dir.glob('*.[Pp][Pp][Tt][Xx]'))

    # Filter out temporary Office files (start with ~$)
    pptx_files = [f for f in pptx_files if not f.name.startswith('~$')]

    if not pptx_files:
        print(f"No PPTX files found in {input_dir}")
        return

    print(f"Found {len(pptx_files)} PPTX file(s) to process")

    for pptx_file in pptx_files:
        try:
            print(f"\nProcessing: {pptx_file.name}")

            # Extract text
            extracted_text = extract_text_from_pptx(pptx_file)

            # Create output markdown file
            base_name = pptx_file.stem
            output_file = output_dir / f"{base_name}.md"

            with open(output_file, "w", encoding="utf-8") as f:
                f.write(extracted_text)

            print(f"  -> {output_file}")

            # Move or keep original
            if keep:
                print(f"  (original kept in input)")
            else:
                processed_file = processed_dir / pptx_file.name
                shutil.move(str(pptx_file), str(processed_file))
                print(f"  -> moved to processed/")

        except Exception as e:
            print(f"  ERROR: {pptx_file.name}: {str(e)}")

    print("\nDone.")


def main():
    parser = argparse.ArgumentParser(description='Extract text from PowerPoint files to markdown')
    parser.add_argument('--dir', '-d', default='.',
                        help='Specify working directory (overrides saved config)')
    parser.add_argument('--setup', action='store_true',
                        help='Run interactive setup to choose working directory')
    parser.add_argument('--keep', '-k', action='store_true',
                        help='Keep original files in input/ instead of moving to processed/')

    args = parser.parse_args()

    if args.setup:
        interactive_setup()
        return

    base_dir = resolve_working_dir(args.dir)
    print(f"Working in: {base_dir}")
    process_pptx_files(base_dir, keep=args.keep)


if __name__ == "__main__":
    main()
